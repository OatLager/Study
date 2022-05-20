## 관련 강의 
# https://www.youtube.com/watch?v=rkABoULImZA

## Excel 관련 설치 =======
# pip install pandas
# pip install openpyxl
## PPT 관련 설치 =========
# pip install python-pptx

# cd Python_example\wors\Excel2PPT
# =================================================================================================

import pandas
from pptx import Presentation

df = pandas.read_excel('Auto_Excel.xlsx',sheet_name='label') # 불러올 excel 파일 이름 및 시트 이름
prs = Presentation('Auto_PPT_Format.pptx')                   # 불러올 ppt 파일 이름

add_slide_layout = prs.slide_layouts[1]         # ppt 슬라이드 마스터에서 추가할 슬라이드 레이아웃 지정.
slide = prs.slides.add_slide(add_slide_layout)  # 지정된 레이아웃 슬라이드 추가 

## ppt 에서 객체 정보(번호) 확인
# placeholder : 객체 
#shapes = slide.shapes

#for shape in shapes:
#    print(str(shape.placeholder_format.idx) + " : " + shape.name)

## 해당 객체에 내용 넣기 : 텍스트 
# Excel에서 내용 가져와 PPT 객체에 채워넣기 
print(df.Label)
slide.placeholders[13].text = df.Label[0]
slide.placeholders[14].text = df.Label[1]
slide.placeholders[15].text = df.Label[2]
slide.placeholders[16].text = df.Label[3]

## PPT 파일 저장
prs.save('Auto_PPT_Format.pptx')